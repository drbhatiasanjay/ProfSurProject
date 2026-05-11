"""
PPTX Generator — assembles Board Deck from topic data dicts.

Uses python-pptx + kaleido (Plotly → PNG → slide).
Each Plotly figure is exported as a high-res PNG and embedded as a picture shape.
Text insight bullets are added as a text box below each chart.
"""

from __future__ import annotations
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from models.board_export import TOPIC_BUILDERS, TOPIC_LABELS

# Slide dimensions — 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Brand colours (hex → int tuples for python-pptx)
_COLOR_HEADER_BG = RGBColor(0x0D, 0x94, 0x88)   # PRIMARY teal
_COLOR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
_COLOR_BODY_FG   = RGBColor(0x1F, 0x29, 0x37)
_COLOR_BULLET_FG = RGBColor(0x37, 0x4A, 0x5F)
_COLOR_ACTION_FG = RGBColor(0x0D, 0x94, 0x88)

# Layout metrics
_HEADER_H   = Inches(0.75)
_FOOTER_H   = Inches(0.30)
_MARGIN_L   = Inches(0.35)
_MARGIN_R   = Inches(0.35)
_CHART_TOP  = _HEADER_H + Inches(0.10)
_CHART_W    = SLIDE_W - _MARGIN_L - _MARGIN_R
_CHART_H    = Inches(4.00)
_TEXT_TOP   = _CHART_TOP + _CHART_H + Inches(0.10)
_TEXT_H     = SLIDE_H - _TEXT_TOP - _FOOTER_H - Inches(0.05)


def _add_header(slide, title: str, company_name: str) -> None:
    """Dark teal header bar with slide title and company name."""
    txBox = slide.shapes.add_textbox(0, 0, SLIDE_W, _HEADER_H)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"{company_name}  ·  {title}"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = _COLOR_HEADER_FG
    # Background fill
    from pptx.util import Pt as _Pt
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = _COLOR_HEADER_BG


def _add_footer(slide, page_num: int, total: int) -> None:
    """Subtle footer with page number."""
    y = SLIDE_H - _FOOTER_H
    txBox = slide.shapes.add_textbox(_MARGIN_L, y, SLIDE_W - _MARGIN_L - _MARGIN_R, _FOOTER_H)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Slide {page_num} of {total}  ·  LifeCycle Leverage Dashboard  ·  ProfSur Analytics"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)


def _add_chart_image(slide, fig, top=None, height=None) -> None:
    """Export a Plotly figure as PNG and embed it on the slide."""
    top = top if top is not None else _CHART_TOP
    height = height if height is not None else _CHART_H
    img_bytes = fig.to_image(format="png", width=1600, height=600, scale=2)
    slide.shapes.add_picture(
        BytesIO(img_bytes),
        _MARGIN_L, top, _CHART_W, height,
    )


def _add_table(slide, df, top=None) -> None:
    """Add a pandas DataFrame as a PPTX table."""
    if df is None or df.empty:
        return
    top = top if top is not None else _TEXT_TOP
    rows, cols = len(df) + 1, len(df.columns)
    col_w = _CHART_W / cols
    tbl = slide.shapes.add_table(rows, cols, _MARGIN_L, top,
                                  _CHART_W, min(_TEXT_H, Inches(0.3 * rows))).table
    tbl.columns[0].width = col_w
    # Header row
    for j, col_name in enumerate(df.columns):
        cell = tbl.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _COLOR_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = _COLOR_HEADER_FG
    # Data rows
    for i, (_, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val) if val is not None else ""
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(8)
            run.font.color.rgb = _COLOR_BODY_FG
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF3, 0xF4, 0xF6)


def _add_insights_text(slide, insights: list[str], actions: list[str]) -> None:
    """Two-column text box: key findings (left) + actions (right)."""
    half_w = (_CHART_W - Inches(0.15)) / 2

    for col_idx, (heading, items, color) in enumerate([
        ("Key Findings", insights, _COLOR_BULLET_FG),
        ("Actions", actions, _COLOR_ACTION_FG),
    ]):
        x = _MARGIN_L + col_idx * (half_w + Inches(0.15))
        txBox = slide.shapes.add_textbox(x, _TEXT_TOP, half_w, _TEXT_H)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Heading
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = _COLOR_HEADER_BG

        # Bullets
        for item in items[:4]:   # max 4 bullets per column
            p = tf.add_paragraph()
            p.space_before = Pt(3)
            run = p.add_run()
            # Strip markdown bold markers for PPTX plain text
            run.text = "• " + item.replace("**", "")
            run.font.size = Pt(8)
            run.font.color.rgb = color


def _cover_slide(prs: Presentation, company_info: dict, selected_count: int) -> None:
    """Branded cover slide."""
    from datetime import date
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Full background block
    bg = slide.shapes.add_textbox(0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _COLOR_HEADER_BG

    # Company name (large)
    txt = slide.shapes.add_textbox(Inches(1), Inches(1.8), SLIDE_W - Inches(2), Inches(1.5))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = company_info.get("name", "Company")
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _COLOR_HEADER_FG

    # Subtitle
    txt2 = slide.shapes.add_textbox(Inches(1), Inches(3.4), SLIDE_W - Inches(2), Inches(0.8))
    tf2 = txt2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        f"Capital Structure Board Deck  ·  {company_info.get('industry', '')}  ·  "
        f"Life Stage: {company_info.get('current_stage', '')}  ·  "
        f"{selected_count} topics selected  ·  {date.today().strftime('%d %b %Y')}"
    )
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0xE0, 0xF2, 0xF1)

    # Footnote
    fn = slide.shapes.add_textbox(Inches(1), Inches(6.6), SLIDE_W - Inches(2), Inches(0.5))
    tf_fn = fn.text_frame
    p_fn = tf_fn.paragraphs[0]
    p_fn.alignment = PP_ALIGN.CENTER
    run_fn = p_fn.add_run()
    run_fn.text = "Generated by LifeCycle Leverage Dashboard — Prof. Surendra Kumar PhD Thesis Platform"
    run_fn.font.size = Pt(9)
    run_fn.font.color.rgb = RGBColor(0xB2, 0xDF, 0xDB)


def build(
    company_df,
    company_info: dict,
    peers_df,
    full_panel,
    stage_summary,
    selected_topics: dict,
) -> bytes:
    """
    Build a PPTX presentation and return as bytes.

    selected_topics: {topic_id: bool} — True = include, False = skip.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    selected_count = sum(v for v in selected_topics.values())
    _cover_slide(prs, company_info, selected_count)

    # Count total content slides for footer
    total_slides = 1 + selected_count  # cover + content
    slide_num = 1

    for topic_id in sorted(selected_topics.keys()):
        if not selected_topics.get(topic_id):
            continue

        builder = TOPIC_BUILDERS.get(topic_id)
        if builder is None:
            continue

        try:
            topic_data = builder(company_df, company_info, peers_df, full_panel, stage_summary)
        except Exception as exc:
            # Never crash the whole deck for one topic — add an error slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_num += 1
            _add_header(slide, TOPIC_LABELS.get(topic_id, f"Topic {topic_id}"), company_info.get("name", ""))
            err_box = slide.shapes.add_textbox(_MARGIN_L, _CHART_TOP, _CHART_W, _CHART_H)
            err_box.text_frame.text = f"Chart unavailable: {exc}"
            _add_footer(slide, slide_num, total_slides)
            continue

        figs    = topic_data.get("figs", [])
        tables  = topic_data.get("tables", [])
        insights = topic_data.get("insights", [])
        actions  = topic_data.get("actions", [])
        title    = topic_data.get("title", TOPIC_LABELS.get(topic_id, ""))

        if figs:
            # First figure: full chart height
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_num += 1
            _add_header(slide, title, company_info.get("name", ""))
            try:
                _add_chart_image(slide, figs[0])
            except Exception:
                pass
            _add_insights_text(slide, insights, actions)
            _add_footer(slide, slide_num, total_slides)

            # Additional figures on their own slides
            for fig in figs[1:]:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide_num += 1
                _add_header(slide, title + " (cont.)", company_info.get("name", ""))
                try:
                    _add_chart_image(slide, fig)
                except Exception:
                    pass
                _add_footer(slide, slide_num, total_slides)

        # Tables on their own slide
        for tbl_df in tables:
            if tbl_df is not None and not tbl_df.empty:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide_num += 1
                _add_header(slide, title + " — Data Table", company_info.get("name", ""))
                try:
                    _add_table(slide, tbl_df.head(12))
                except Exception:
                    pass
                _add_footer(slide, slide_num, total_slides)

        # If topic has no figs but has insights (e.g. text-only topic)
        if not figs and (insights or actions):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_num += 1
            _add_header(slide, title, company_info.get("name", ""))
            _add_insights_text(slide, insights, actions)
            _add_footer(slide, slide_num, total_slides)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
